# document_processor.py - Document Processing and Chunking
import asyncio
import logging
from typing import List, Dict, Any, Optional
from urllib.parse import urlparse
import aiohttp
import tempfile
import os
from pathlib import Path
import hashlib
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
import PyPDF2
from docx import Document as DocxDocument
import tiktoken
from config import settings
import io
import zipfile
import re
import xml.etree.ElementTree as ET
import base64
import openai

logger = logging.getLogger(__name__)

class DocumentChunk:
    def __init__(self, content: str, metadata: Dict[str, Any]):
        self.content = content
        self.metadata = metadata
        self.page_number = metadata.get('page_number')
        self.chunk_id = metadata.get('chunk_id')
        self.source = metadata.get('source')

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,  # Increased for better context preservation
            chunk_overlap=200,  # Increased for better context continuity
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        self.session = None
        self.cache = {}  # Cache for processed documents
        self.supported_extensions = {".pdf", ".docx", ".txt", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
        # Classification sets for unsupported formats
        self.image_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".svg"}
        self.archive_extensions = {".zip", ".tar", ".gz", ".tgz", ".rar", ".7z", ".xz", ".bz2"}
        self.binary_extensions = {".bin", ".exe", ".dll", ".so", ".dylib", ".pkg", ".deb", ".rpm"}
        self.openai_client = None

    async def initialize(self):
        """Initialize aiohttp session for document downloads"""
        if not self.session:
            # Unlimited timeout to support very large documents
            self.session = aiohttp.ClientSession(
                timeout=aiohttp.ClientTimeout(total=None)
            )
        # Initialize OpenAI client for OCR if available
        if self.openai_client is None and settings.OPENAI_API_KEY and settings.OPENAI_API_KEY != "your_openai_api_key_here":
            try:
                self.openai_client = openai.AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
            except Exception:
                self.openai_client = None

    async def process_document(self, document_url: str) -> List[Document]:
        """Process document with ultra-fast optimization"""
        try:
            # Validate supported document type early
            ext = self._extract_extension_from_source(document_url)
            if ext not in self.supported_extensions:
                logger.warning(f"Unsupported document format '{ext}' for URL: {document_url}")
                return []
            # Check cache first
            cache_key = hashlib.md5(document_url.encode()).hexdigest()
            if cache_key in self.cache:
                return self.cache[cache_key]

            # Download document
            content = await self._download_document(document_url)
            if not content:
                return []

            # Extract text based on file type
            text = await self._extract_text(content, document_url)
            if not text:
                return []

            # Split text into chunks
            chunks = await self._split_text(text, document_url)
            
            # Cache the result
            self.cache[cache_key] = chunks
            return chunks

        except Exception as e:
            logger.error(f"Error processing document {document_url}: {str(e)}")
            return []

    async def _download_document(self, url: str) -> Optional[bytes]:
        """Download document with ultra-fast optimization"""
        try:
            if not self.session:
                await self.initialize()

            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; HackrxDocFetcher/1.0)",
                "Accept": "*/*"
            }
            async with self.session.get(
                url,
                headers=headers,
                timeout=aiohttp.ClientTimeout(total=None)
            ) as response:
                if response.status == 200:
                    return await response.read()
                else:
                    logger.error(
                        f"Failed to download document from {url}: HTTP {response.status}"
                    )
                    return None

        except Exception as e:
            logger.error(f"Error downloading document from {url}: {str(e)}")
            return None

    async def _extract_text(self, content: bytes, url: str) -> Optional[str]:
        """Extract text from document with ultra-fast optimization"""
        try:
            # Determine extension robustly (handles URLs with query params)
            file_extension = self._extract_extension_from_source(url)
            
            if file_extension == '.pdf':
                return await self._extract_pdf_text(content)
            elif file_extension in ['.docx', '.doc']:
                # Primary extraction via python-docx
                text = await self._extract_docx_text(content)
                # Fallback to XML unzip if empty or very short
                if not text or len(text.strip()) < 10:
                    fallback_text = self._extract_docx_text_fallback(content)
                    return fallback_text
                return text
            elif file_extension == '.xlsx':
                return await self._extract_xlsx_text(content)
            elif file_extension in ['.pptx', '.ppt']:
                return await self._extract_pptx_text(content)
            elif file_extension in self.image_extensions:
                # OCR via OpenAI Vision (gpt-4o)
                return await self._extract_image_text_openai(content)
            else:
                # Assume it's plain text
                return content.decode('utf-8', errors='ignore')

        except Exception as e:
            logger.error(f"Error extracting text from document: {str(e)}")
            return None

    async def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF with ultra-fast optimization"""
        try:
            text = ""
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            
            # Process all pages for complete coverage
            max_pages = len(pdf_reader.pages)
            for page_num in range(max_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
            
            return text

        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return ""

    async def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX with ultra-fast optimization"""
        try:
            doc = DocxDocument(io.BytesIO(content))
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            return text

        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            return ""

    def _extract_docx_text_fallback(self, content: bytes) -> str:
        """Fallback DOCX extraction by parsing word/document.xml and concatenating w:t elements."""
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if 'word/document.xml' not in zf.namelist():
                    return ""
                with zf.open('word/document.xml') as f:
                    xml_text = f.read().decode('utf-8', errors='ignore')
                    # Parse XML and extract all text nodes
                    try:
                        # Register namespaces if present
                        # Extract text from all w:t tags, inserting newlines for paragraphs
                        root = ET.fromstring(xml_text)
                        namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                        texts: List[str] = []
                        for para in root.findall('.//w:p', namespaces):
                            parts: List[str] = []
                            for node in para.findall('.//w:t', namespaces):
                                if node.text:
                                    parts.append(node.text)
                            if parts:
                                texts.append(' '.join(parts).strip())
                        plain_text = '\n'.join(texts)
                        # Normalize whitespace
                        plain_text = re.sub(r"\s+", " ", plain_text)
                        plain_text = re.sub(r"\s*\n\s*", "\n", plain_text).strip()
                        return plain_text
                    except Exception:
                        # As last resort, strip tags
                        stripped = re.sub(r'<[^>]+>', ' ', xml_text)
                        stripped = re.sub(r"\s+", " ", stripped).strip()
                        return stripped
        except Exception as e:
            logger.error(f"Error in DOCX fallback extraction: {str(e)}")
            return ""

    async def _extract_image_text_openai(self, content: bytes) -> str:
        """Use OpenAI vision to extract raw text from image bytes."""
        try:
            if not self.openai_client:
                logger.warning("OpenAI client not available for OCR - returning empty text")
                return ""
            # Encode image as base64 data URL
            b64 = base64.b64encode(content).decode('utf-8')
            data_url = f"data:image/png;base64,{b64}"
            response = await self.openai_client.chat.completions.create(
                model=settings.OPENAI_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract only the raw text from this image in natural reading order. No commentary."},
                            {"type": "image_url", "image_url": {"url": data_url}}
                        ]
                    }
                ],
                temperature=0.0,
                max_tokens=2000
            )
            text = response.choices[0].message.content or ""
            return text.strip()
        except Exception as e:
            logger.error(f"Error during image OCR with OpenAI: {str(e)}")
            return ""

    async def _extract_xlsx_text(self, content: bytes) -> str:
        """Extract text from XLSX spreadsheets by reading all sheets and rows."""
        try:
            # Lazy import to avoid hard dependency if not installed
            import openpyxl  # type: ignore
        except Exception as e:
            logger.error(f"openpyxl not available for XLSX parsing: {str(e)}")
            return ""

        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
            parts: List[str] = []
            for sheet in workbook.worksheets:
                parts.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    # Convert row to string, skipping completely empty rows
                    values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in values):
                        parts.append(" | ".join(values))
                parts.append("")
            text = "\n".join(parts)
            return text
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {str(e)}")
            return ""

    async def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PPTX slides including shapes, tables, and notes."""
        try:
            # Lazy import to avoid hard dependency if not installed
            from pptx import Presentation  # type: ignore
        except Exception as e:
            logger.error(f"python-pptx not available for PPTX parsing: {str(e)}")
            return ""

        try:
            prs = Presentation(io.BytesIO(content))
            parts: List[str] = []
            for idx, slide in enumerate(prs.slides, start=1):
                parts.append(f"Slide {idx}:")
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Text frames
                    if getattr(shape, "has_text_frame", False):
                        text_runs: List[str] = []
                        for paragraph in shape.text_frame.paragraphs:
                            run_text = "".join(run.text for run in paragraph.runs) or paragraph.text
                            if run_text and run_text.strip():
                                text_runs.append(run_text.strip())
                        if text_runs:
                            parts.append(" ".join(text_runs))
                    # Tables
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() if cell.text else "" for cell in row.cells]
                            if any(cells):
                                parts.append(" | ".join(cells))
                # Notes (if present)
                try:
                    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = []
                        for p in slide.notes_slide.notes_text_frame.paragraphs:
                            t = "".join(r.text for r in p.runs) or p.text
                            if t and t.strip():
                                notes_text.append(t.strip())
                        if notes_text:
                            parts.append("Notes: " + " ".join(notes_text))
                except Exception:
                    pass
                parts.append("")
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {str(e)}")
            return ""

    def _extract_extension_from_source(self, source: str) -> str:
        """Extract a safe file extension from a URL or file path, ignoring query params."""
        try:
            parsed = urlparse(source)
            path = parsed.path if parsed.scheme else source
            ext = Path(path).suffix.lower()
            # Normalize to common known types; keep it short and safe
            if ext and len(ext) <= 10:
                return ext
            # Fallback: best-effort guess
            for candidate in [".pdf", ".docx", ".doc", ".txt", ".zip"]:
                if candidate in path.lower():
                    return candidate
            return ""
        except Exception:
            return ""

    def is_supported_source(self, source: str) -> bool:
        """Check if the source URL/path has a supported extension."""
        ext = self._extract_extension_from_source(source)
        return ext in self.supported_extensions

    def classify_source(self, source: str) -> Dict[str, str]:
        """Classify the source by extension into a human-friendly category.

        Returns: { 'extension': str, 'category': str }
        Categories: pdf, docx, text, image, archive, binary, unknown
        """
        ext = self._extract_extension_from_source(source)
        ext_lower = ext.lower() if ext else ""
        if ext_lower in {".pdf"}:
            return {"extension": ext_lower, "category": "pdf"}
        if ext_lower in {".docx"}:
            return {"extension": ext_lower, "category": "docx"}
        if ext_lower in {".txt"}:
            return {"extension": ext_lower, "category": "text"}
        if ext_lower in {".pptx", ".ppt"}:
            return {"extension": ext_lower, "category": "pptx"}
        if ext_lower in self.image_extensions:
            return {"extension": ext_lower, "category": "image"}
        if ext_lower in self.archive_extensions:
            return {"extension": ext_lower, "category": "archive"}
        if ext_lower in self.binary_extensions:
            return {"extension": ext_lower, "category": "binary"}
        return {"extension": ext_lower or "", "category": "unknown"}

    async def _split_text(self, text: str, source: str) -> List[Document]:
        """Split text into chunks with ultra-fast optimization"""
        try:
            if not text.strip():
                return []

            # Use ultra-fast text splitting
            chunks = self.text_splitter.split_text(text)
            
            # Create Document objects with metadata
            documents = []
            safe_ext = self._extract_extension_from_source(source)
            for i, chunk in enumerate(chunks):
                doc = Document(
                    page_content=chunk,
                    metadata={
                        'source': source,
                        'chunk_id': f"{hashlib.md5(source.encode()).hexdigest()}_{i}",
                        'chunk_index': i,
                        'total_chunks': len(chunks),
                        'document_type': safe_ext,
                        'token_count': len(chunk.split())
                    }
                )
                documents.append(doc)

            return documents

        except Exception as e:
            logger.error(f"Error splitting text: {str(e)}")
            return []

    async def cleanup(self):
        """Cleanup resources"""
        if self.session:
            await self.session.close()
    